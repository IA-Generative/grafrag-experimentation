from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


TEMPLATE_PATH = Path(
    "/Users/etiquet/Documents/GitHub/PIP-Organization-Chart-From-Grist/data/template.ppt.pptx"
)
OUT_DIR = Path("presentations")
PPTX_PATH = OUT_DIR / "demo_grafrag_anef_10_15min_template.pptx"
MD_PATH = OUT_DIR / "demo_grafrag_anef_10_15min_template.md"


SLIDES: list[dict[str, object]] = [
    {
        "title": "MirAI - du graphe documentaire a l'assistant reglementaire",
        "subtitle": "Demo 10-15 min | grafrag-experimentation + anef-knowledge-assistant",
        "notes": [
            "Ouvrir sur l'idee d'une meme experience utilisateur pour plusieurs intelligences specialisees.",
            "Le message a faire passer: ce n'est pas deux demos cote a cote, c'est une plateforme d'assistants.",
        ],
    },
    {
        "title": "POURQUOI CES EXPERIMENTATIONS COMPTENT",
        "boxes": [
            "Intention\nProuver qu'un chat peut devenir un assistant fiable, sourcable et actionnable.\nTester une UX unique pour retrieval, graphe, citations et outils metier.",
            "Ce que montre grafrag\nGraphRAG derriere Open WebUI.\nCorpus Manager, viewer interactif, fallback, cycle de vie des corpus.",
            "Ce que montre ANEF\nMoteur d'eligibilite explicable.\nCitations CESEDA cliquables, legal viewer, surcouche browser legere.",
        ],
        "notes": [
            "Cette slide doit sonner marketing: on parle valeur, confiance et demonstrabilite.",
            "Insister sur la complementarite entre exploration documentaire et decision assistee.",
        ],
    },
    {
        "title": "CE QUE PEUT FAIRE GRAFRAG-EXPERIMENTATION",
        "boxes": [
            "Valeur demo\nQuestionner un corpus, obtenir une reponse sourcee, puis naviguer dans le graphe et la chronologie.",
            "Capacites clefs\nBridge FastAPI + pipelines Open WebUI.\nGraph viewer.\nFallback documentaire si l'index complet manque.\nMode multi-corpus avec sync, index, publish et ACL.",
        ],
        "notes": [
            "Montrer que le repo ne sert pas seulement a lancer GraphRAG, mais a operer un workflow de corpus.",
            "Dire explicitement que la resilience fait partie de la proposition de valeur.",
        ],
    },
    {
        "title": "CE QUE PEUT FAIRE ANEF-KNOWLEDGE-ASSISTANT",
        "boxes": [
            "Valeur demo\nTransformer un corpus reglementaire et un Excel metier en assistant d'eligibilite explicable.",
            "Capacites clefs\nAPI FastAPI orientee usages.\nPieces, conditions, wizard, FAQ, legal search.\nViewer CESEDA local.\nReponses groundees avec citations cliquables et preview au survol.",
        ],
        "notes": [
            "Positionner ANEF comme une preuve d'usage vertical a forte valeur metier.",
            "Le point fort n'est pas seulement la recherche, mais l'explication et la vigilance.",
        ],
    },
    {
        "title": "POURQUOI LES DEUX REPOS ENSEMBLE SONT INTERESSANTS",
        "boxes": [
            "Experience commune\nOpen WebUI, SSO Keycloak, aliases MirAI, pipelines et surface conversationnelle partages.",
            "Architecture et garde-fous\nRepos separes, responsabilites claires, redeploiement partage, reprovisioning des modeles et du loader.js.",
            "Synthese\nUne plateforme modulaire ou chaque assistant specialise peut evoluer sans casser l'experience commune.",
        ],
        "notes": [
            "C'est la slide d'architecture racontee en langage produit.",
            "Expliquer que le choix de deux repos permet d'accelerer sans tout coupler.",
        ],
    },
    {
        "title": "FINALITES ET AMBITION DE LA DEMO",
        "boxes": [
            "Finalites\nMontrer une experience complete: poser une question, voir les sources, naviguer, comprendre et agir.\nMontrer aussi que le systeme reste presentable meme quand tout n'est pas parfait.",
            "Ambition\nProuver qu'Open WebUI peut devenir une couche d'experience commune pour plusieurs moteurs IA specialises, du graphe documentaire au conseil reglementaire.",
        ],
        "notes": [
            "Ne pas survendre l'instantane: l'ambition est la plateforme demonstrable et industrialisable.",
            "Cette slide sert de transition avant la demo live.",
        ],
    },
    {
        "title": "DEROULE DE DEMO ET TESTS",
        "boxes": [
            "1. Ouvrir Open WebUI et presenter les modeles MirAI.\n2. Prompt GraphRAG: chronologie de la guerre de Cent Ans et batailles pivots.\n3. Ouvrir le graph viewer pour montrer noeuds, relations et lecture chronologique.\n4. Montrer le workflow Corpus Manager: sync > index > publish, ou l'expliquer si l'index tourne encore.\n5. Prompt ANEF: CST salarie - L.421-1, pieces, conditions et vigilance.\n6. Survoler une citation CESEDA, ouvrir la modale, puis conclure sur la valeur plateforme.",
        ],
        "notes": [
            "Si le mode multi-corpus est actif, prefixer le prompt GraphRAG avec [[corpus:<id>]].",
            "Si l'index GraphRAG n'est pas fini, assumer le fallback comme une preuve de resilience.",
            "Terminer en rappelant les limites: indexation longue, aliases rejouables, validation humaine indispensable sur ANEF.",
        ],
    },
    {
        "title": "FEEDBACK, STRATEGIE DE PROMPT ET ENSEIGNEMENTS CLES",
        "boxes": [
            "Feedback a rechercher\nConfiance grace aux citations.\nNavigation entre chat, graphe et viewer juridique.\nSensation de rapidite meme quand l'index n'est pas parfait.",
            "Strategie de prompt\nGraphRAG: demander chronologie, acteurs, comparaisons, contexte.\nMulti-corpus: prefixer [[corpus:<id>]] si necessaire.\nANEF: demander pieces, conditions, base legale et points de vigilance dans une meme requete.",
            "Enseignements cles\nLa resilience compte plus que la perfection.\nLe fallback et les timeouts protegent l'experience.\nUn loader.js leger et des scripts de reprovisioning valent mieux qu'un fork frontal lourd.\nLa validation humaine reste indispensable sur le reglementaire.",
        ],
        "notes": [
            "Cette slide permet de finir avec une lecture mature du systeme: valeur, bon usage, et limites assumees.",
            "Les lessons learned viennent surtout du README global: fallback, timeouts, reprovisioning, separation des repos, et webui.db volatil.",
        ],
    },
]

MERMAID_SECTIONS: list[dict[str, str]] = [
    {
        "title": "Architecture GraphRAG simplifiee",
        "intro": "Schema directement inspire du README global pour expliquer le trajet d'une question jusqu'a la reponse sourcee.",
        "code": """flowchart TD
    U[Utilisateur] --> O[Open WebUI]
    O --> P[Pipeline Open WebUI]
    P --> B[Bridge FastAPI GraphRAG]
    B --> G[Microsoft GraphRAG]
    G --> L[LLM OpenAI-compatible]
    B --> C[(Corpus local ou version publiee)]""",
    },
    {
        "title": "Topologie partagee des deux repositories",
        "intro": "Vue d'architecture mutualisee quand `grafrag-experimentation` et `anef-knowledge-assistant` ciblent la meme instance Open WebUI.",
        "code": """flowchart LR
    U[Browser utilisateur] --> OWF[Open WebUI frontend]
    OWF --> OWB[Open WebUI backend]
    OWB --> P[pipelines service]
    OWF -. OIDC .-> KC[Keycloak]

    subgraph G[grafrag-experimentation]
        BR[GraphRAG bridge]
        CM[Corpus Manager API]
        CW[corpus-worker]
        WS[(Workspaces GraphRAG versionnes)]
        SRC[Sources synchronisees]
    end

    subgraph A[anef-knowledge-assistant]
        AP[ANEF API]
        AR[Pipeline reglementaire]
        LJ[loader.js CESEDA]
    end

    P --> BR
    P --> AR
    BR --> WS
    CM --> CW
    CW --> SRC
    OWF --> LJ
    LJ --> AP""",
    },
    {
        "title": "Fonctionnement de la requete et des garde-fous",
        "intro": "Le point important a expliciter pendant la demo: le systeme privilegie la continuite de service et la lisibilite des reponses.",
        "code": """flowchart TD
    Q[Question utilisateur] --> R{Type de demande}
    R -->|Documentaire| GP[Pipeline GraphRAG]
    R -->|Reglementaire| AP[Pipeline ANEF]

    GP --> T{Index disponible et budget temps suffisant ?}
    T -->|Oui| CLI[graphrag query]
    T -->|Non| FB[Fallback documentaire]
    CLI --> SYN[Reponse synthetisee + sources + lien graphe]
    FB --> SYN

    AP --> TOOLS[eligibility-check / legal-search / search-title]
    TOOLS --> LEGAL[Reponse groundee avec citations CESEDA]
    LEGAL --> UX[Preview au survol + modal au clic]
    SYN --> UX""",
    },
]


def build_markdown() -> str:
    lines: list[str] = [
        "# Demo MirAI 10-15 min - version template",
        "",
        f"Template source: `{TEMPLATE_PATH}`",
        "",
    ]
    for index, slide in enumerate(SLIDES, start=1):
        lines.append(f"## Slide {index} - {slide['title']}")
        lines.append("")
        subtitle = slide.get("subtitle")
        if isinstance(subtitle, str):
            lines.append(f"- {subtitle}")
            lines.append("")
        boxes = slide.get("boxes")
        if isinstance(boxes, list):
            for box_index, box in enumerate(boxes, start=1):
                lines.append(f"Bloc {box_index}:")
                for raw_line in str(box).splitlines():
                    if raw_line.strip():
                        lines.append(f"- {raw_line.strip()}")
                lines.append("")
        notes = slide.get("notes")
        if isinstance(notes, list):
            lines.append("Notes orateur:")
            for note in notes:
                lines.append(f"- {note}")
            lines.append("")
        if index == 5:
            lines.append("## Schemas Mermaid d'appui")
            lines.append("")
            for section in MERMAID_SECTIONS:
                lines.append(f"### {section['title']}")
                lines.append("")
                lines.append(section["intro"])
                lines.append("")
                lines.append("```mermaid")
                lines.append(section["code"])
                lines.append("```")
                lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def set_shape_text(shape, text: str) -> None:
    shape.text = text


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, raw_line in enumerate(text.splitlines()):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = raw_line
        if index == 0:
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)
        else:
            paragraph.font.size = Pt(14)


def build_presentation() -> Presentation:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    prs = Presentation(str(TEMPLATE_PATH))

    # Slide 1
    set_shape_text(prs.slides[0].shapes[0], str(SLIDES[0]["title"]))
    set_shape_text(prs.slides[0].shapes[1], str(SLIDES[0]["subtitle"]))

    # Slide 2
    slide = prs.slides[1]
    set_shape_text(slide.shapes[0], str(SLIDES[1]["title"]))
    for idx, box in enumerate(SLIDES[1]["boxes"], start=1):
        set_shape_text(slide.shapes[idx], str(box))

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide.shapes[0], str(SLIDES[2]["title"]))
    for idx, box in enumerate(SLIDES[2]["boxes"], start=1):
        set_shape_text(slide.shapes[idx], str(box))

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide.shapes[0], str(SLIDES[3]["title"]))
    for idx, box in enumerate(SLIDES[3]["boxes"], start=1):
        set_shape_text(slide.shapes[idx], str(box))

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide.shapes[0], str(SLIDES[4]["title"]))
    for idx, box in enumerate(SLIDES[4]["boxes"], start=1):
        set_shape_text(slide.shapes[idx], str(box))

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide.shapes[0], str(SLIDES[5]["title"]))
    for idx, box in enumerate(SLIDES[5]["boxes"], start=1):
        set_shape_text(slide.shapes[idx], str(box))

    # Slide 7
    slide = prs.slides[6]
    set_shape_text(slide.shapes[0], str(SLIDES[6]["title"]))
    set_shape_text(slide.shapes[1], str(SLIDES[6]["boxes"][0]))

    # Slide 8 - custom slide based on template layout
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    set_shape_text(slide.shapes[0], str(SLIDES[7]["title"]))
    add_textbox(slide, 0.9, 1.55, 3.8, 4.5, str(SLIDES[7]["boxes"][0]))
    add_textbox(slide, 4.8, 1.55, 3.8, 4.5, str(SLIDES[7]["boxes"][1]))
    add_textbox(slide, 8.7, 1.55, 3.8, 4.8, str(SLIDES[7]["boxes"][2]))

    return prs


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    MD_PATH.write_text(build_markdown(), encoding="utf-8")
    prs = build_presentation()
    prs.save(PPTX_PATH)
    print(f"Wrote {MD_PATH}")
    print(f"Wrote {PPTX_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
