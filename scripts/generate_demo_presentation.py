from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("presentations")
PPTX_PATH = OUT_DIR / "demo_grafrag_anef_10_15min.pptx"
MD_PATH = OUT_DIR / "demo_grafrag_anef_10_15min.md"

NAVY = RGBColor(14, 37, 66)
TEAL = RGBColor(25, 125, 122)
SAND = RGBColor(242, 236, 226)
SLATE = RGBColor(73, 89, 106)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(32, 41, 51)
MUTED = RGBColor(106, 117, 129)


SLIDES: list[dict[str, object]] = [
    {
        "kind": "title",
        "title": "MirAI: du graphe documentaire a l'assistant reglementaire",
        "subtitle": [
            "Support de demonstration 10-15 min",
            "grafrag-experimentation + anef-knowledge-assistant",
            "Une meme experience Open WebUI pour deux usages a forte valeur",
        ],
        "notes": [
            "Ouvrir en rappelant que l'objectif n'est pas de montrer deux prototypes isoles.",
            "Le message cle: une meme surface conversationnelle peut accueillir un assistant GraphRAG documentaire et un assistant metier reglementaire.",
            "Positionner la demo comme une experimentation produit, technique et d'industrialisation.",
        ],
    },
    {
        "kind": "bullets",
        "title": "L'intention de l'experimentation",
        "bullets": [
            "Tester une experience IA plus fiable qu'un simple chat generaliste.",
            "Combiner retrieval, graphe, citations et outils metier dans une meme interface.",
            "Prouver qu'on peut faire cohabiter usages generiques et expertise verticale sans forker Open WebUI.",
            "Valider un chemin reproductible du laptop au cluster Kubernetes.",
        ],
        "notes": [
            "Cette slide doit sonner produit: on cherche de la confiance, de l'explicabilite et une experience unifiee.",
            "Insister sur le fait que le repo sert autant a apprendre ce qui marche qu'a derisquer un futur produit.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Ce que demontre grafrag-experimentation",
        "bullets": [
            "GraphRAG branche derriere Open WebUI via un bridge FastAPI et des pipelines dedies.",
            "Visualisation du corpus dans un graph viewer interactif, avec lecture documentaire et chronologique.",
            "Mode resilient: fallback documentaire si l'index complet n'est pas pret.",
            "Mode multi-corpus: Corpus Manager, versions publiees, ACL, worker d'indexation et sources synchronisees.",
        ],
        "notes": [
            "Faire passer l'idee que ce repo ne sert pas seulement a indexer un corpus, mais a operer un cycle de vie de corpus.",
            "Mentionner la valeur demo: meme sans index complet, le viewer et les reponses restent utilisables.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Ce que demontre anef-knowledge-assistant",
        "bullets": [
            "Transformation d'un Excel metier et du CESEDA en moteur d'eligibilite explicable.",
            "API FastAPI orientee cas d'usage: eligibility, pieces, conditions, wizard, FAQ, legal search.",
            "Reponses groundees avec citations CESEDA cliquables et viewer juridique local.",
            "Surcouche browser legere dans Open WebUI: preview au survol, modal au clic, sans fork frontal lourd.",
        ],
        "notes": [
            "Positionner ANEF comme la preuve qu'on peut aller au-dela de la recherche documentaire vers la decision assistee.",
            "Le point fort marketing est l'explication reglementaire, pas seulement la recherche.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Pourquoi les deux repos ensemble sont interessants",
        "bullets": [
            "Une meme entree utilisateur: Open WebUI, SSO Keycloak, aliases MirAI, pipelines partages.",
            "Deux moteurs complementaires: exploration documentaire d'un cote, expertise reglementaire de l'autre.",
            "Un redeploiement partage capable de reprovisionner modeles, grants, loader.js et integrations.",
            "Une architecture modulaire: chaque domaine evolue sans casser l'autre.",
        ],
        "notes": [
            "Ici il faut raconter le choix d'architecture: separer les responsabilites, mutualiser l'experience.",
            "Insister sur la valeur d'une plateforme d'assistants plutot qu'un monolithe.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Scenario de demo en 10-15 min",
        "bullets": [
            "1. Ouvrir Open WebUI et presenter les modeles MirAI disponibles.",
            "2. Montrer une question historique sur le corpus medieval et la reponse GraphRAG avec sources.",
            "3. Ouvrir le graph viewer pour visualiser les noeuds, relations et la chronologie.",
            "4. Montrer le Corpus Manager ou au minimum expliquer le workflow sync > index > publish.",
            "5. Basculer sur ANEF pour un cas reglementaire concret et afficher les citations CESEDA enrichies.",
        ],
        "notes": [
            "Garder un rythme dynamique: 2 a 3 min GraphRAG, 2 min viewer, 2 min cycle de corpus, 3 min ANEF, 1 min conclusion.",
            "Ne pas faire toute l'admin en live si l'indexation est longue: montrer les etapes, puis basculer sur un resultat prepare.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Deroule de test recommande",
        "bullets": [
            "Prompt GraphRAG: Donne-moi une chronologie synthetique de la guerre de Cent Ans avec les batailles pivots.",
            "Prompt GraphRAG: Quels acteurs relient Crecy, Poitiers, Azincourt et le traite de Troyes ?",
            "Prompt ANEF: Pour un CST salarie - L. 421-1, quelles pieces et quels points de vigilance ?",
            "Interaction ANEF: survoler un lien CESEDA, ouvrir la modale, verifier la citation et le permalink.",
        ],
        "notes": [
            "Si le mode multi-corpus est actif, prefixer le prompt GraphRAG avec [[corpus:<id>]].",
            "Si l'index n'est pas termine, utiliser le mode fallback et assumer que la demonstration porte aussi sur la resilience.",
            "Pour ANEF, viser une reponse qui montre pieces, conditions, citations, vigilance et revue humaine.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Messages techniques a faire passer",
        "bullets": [
            "Le bridge GraphRAG impose des timeouts et un fallback pour garder l'UI reactive.",
            "Le viewer reste disponible meme sans artefacts complets grace au document-map fallback.",
            "Le moteur ANEF garde les citations et l'explicabilite comme premier livrable, pas comme post-traitement cosmetique.",
            "Le loader.js ANEF est persiste declarativement via ConfigMap, contrairement aux aliases Open WebUI rejoues par script.",
        ],
        "notes": [
            "Cette slide sert a rassurer un public technique ou sponsor: la demo n'est pas un bricolage ponctuel.",
            "Mettre en avant les choix d'operabilite: reprovisioning, resilience, separation des concerns.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Valeur metier et valeur de plateforme",
        "bullets": [
            "Pour l'utilisateur: reponses plus actionnables, plus sourcables, plus navigables.",
            "Pour l'equipe produit: experimentation rapide sur plusieurs UX sans reecrire toute la stack.",
            "Pour l'IT: local Docker et Kubernetes partagent une logique d'integration coherente.",
            "Pour le sponsor: preuve qu'un assistant de confiance peut mixer retrieval, graphe, moteur metier et UX augmentee.",
        ],
        "notes": [
            "C'est la slide la plus marketing. Parler resultat et non seulement composants.",
            "Le message final: on construit une base pour des assistants specialises, pas une simple demo technique.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Limites assumees a verbaliser",
        "bullets": [
            "GraphRAG peut etre long a indexer sur des corpus reels: la demo doit montrer le workflow, pas promettre l'instantane.",
            "Les aliases Open WebUI restent rejouables apres rollout si webui.db est volatil.",
            "Le mode multi-corpus existe, mais n'est pas encore une plateforme de gouvernance a grande echelle.",
            "ANEF assiste la decision et met en avant les zones d'incertitude; il ne remplace pas une validation humaine.",
        ],
        "notes": [
            "Bien formuler les limites augmente la credibilite de la demo.",
            "Le bon ton: nous savons ou sont les bords du systeme et nous avons des garde-fous.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Conclusion et ouverture",
        "bullets": [
            "Ce que nous testons: une plateforme d'assistants specialises, demonstrable et industrialisable.",
            "Ce que nous prouvons: Open WebUI peut devenir une couche d'experience commune pour plusieurs moteurs IA.",
            "Ce que nous ouvrons: nouveaux corpus, nouveaux moteurs metier, nouveaux workflows de publication et de controle.",
            "Prochaine etape possible: choisir 1 ou 2 cas d'usage prioritaires et durcir la boucle de validation.",
        ],
        "notes": [
            "Fermer en revenant a la these initiale: une meme experience, plusieurs intelligences specialisees.",
            "Inviter la discussion sur les cas d'usage a prioriser apres la demo.",
        ],
    },
]


def build_markdown() -> str:
    lines: list[str] = [
        "# Demo MirAI 10-15 min",
        "",
        "Support de presentation pour `grafrag-experimentation` et `anef-knowledge-assistant`.",
        "",
    ]
    for index, slide in enumerate(SLIDES, start=1):
        title = str(slide["title"])
        lines.append(f"## Slide {index} - {title}")
        lines.append("")
        subtitle = slide.get("subtitle")
        if isinstance(subtitle, list):
            for item in subtitle:
                lines.append(f"- {item}")
            lines.append("")
        bullets = slide.get("bullets")
        if isinstance(bullets, list):
            for item in bullets:
                lines.append(f"- {item}")
            lines.append("")
        notes = slide.get("notes")
        if isinstance(notes, list):
            lines.append("Notes orateur:")
            for item in notes:
                lines.append(f"- {item}")
            lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.25))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED
    paragraph.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs: Presentation, slide_data: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(13.333),
        Inches(0.35),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = str(slide_data["title"])
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.85), Inches(3.2), Inches(9.7), Inches(2.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.margin_left = 0
    subtitle_frame.margin_right = 0
    subtitle_frame.margin_top = 0
    subtitle_frame.margin_bottom = 0
    subtitles = slide_data.get("subtitle", [])
    if isinstance(subtitles, list):
        for index, item in enumerate(subtitles):
            paragraph = subtitle_frame.paragraphs[0] if index == 0 else subtitle_frame.add_paragraph()
            paragraph.text = str(item)
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = SAND
            paragraph.space_after = Pt(6)

    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(5.8),
        Inches(4.1),
        Inches(0.5),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(38, 63, 97)
    chip.line.fill.background()
    chip_text = chip.text_frame
    chip_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    chip_paragraph = chip_text.paragraphs[0]
    chip_paragraph.alignment = PP_ALIGN.CENTER
    chip_run = chip_paragraph.add_run()
    chip_run.text = "Narration produit + preuve technique"
    chip_run.font.size = Pt(14)
    chip_run.font.color.rgb = WHITE

    add_footer(slide, "MirAI demo deck | 10-15 min")


def add_bullet_slide(prs: Presentation, slide_data: dict[str, object], slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    header_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(13.333),
        Inches(0.4),
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = NAVY
    header_bar.line.fill.background()

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(1.0),
        Inches(0.18),
        Inches(5.2),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TEAL
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.1), Inches(0.8))
    title_frame = title_box.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = str(slide_data["title"])
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = NAVY

    body_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.2), Inches(4.9))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.margin_left = 0
    body_frame.margin_right = 0
    bullets = slide_data.get("bullets", [])
    if isinstance(bullets, list):
        for index, item in enumerate(bullets):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = str(item)
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = DARK
            paragraph.space_after = Pt(11)

    note_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.8),
        Inches(6.35),
        Inches(3.85),
        Inches(0.42),
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = SAND
    note_box.line.color.rgb = RGBColor(221, 211, 197)
    note_frame = note_box.text_frame
    note_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    note_paragraph = note_frame.paragraphs[0]
    note_paragraph.alignment = PP_ALIGN.CENTER
    note_run = note_paragraph.add_run()
    note_run.text = "Demo narrative + support speaker notes"
    note_run.font.size = Pt(11)
    note_run.font.color.rgb = SLATE

    add_footer(slide, f"Slide {slide_number} | MirAI demo deck")


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, slide_data in enumerate(SLIDES, start=1):
        if slide_data["kind"] == "title":
            add_title_slide(prs, slide_data)
        else:
            add_bullet_slide(prs, slide_data, index)
    return prs


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    MD_PATH.write_text(build_markdown(), encoding="utf-8")
    prs = build_presentation()
    prs.save(PPTX_PATH)
    print(f"Wrote {MD_PATH}")
    print(f"Wrote {PPTX_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
